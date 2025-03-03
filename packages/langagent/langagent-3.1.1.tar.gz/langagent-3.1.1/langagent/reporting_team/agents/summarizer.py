# LangAgent/reporting_team/agents/summarizer.py


# * LIBRARIES

import os
import yaml
from langchain_community.document_loaders import PyPDFLoader
from langchain.chains.llm import LLMChain
from langchain_core.prompts import PromptTemplate
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
from langchain_community.chat_models import ChatOllama
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from typing import List, TypedDict
from langgraph.graph import END, StateGraph
from langchain.schema import Document
from langchain_groq import ChatGroq


def create_summarizer(llm):
    
    # Function to detect the input type and load a document
    def load_document(input_data):
        """
        Load a document of specified type (PDF, PPTX, DOCX, or text) and convert it into a list of Document objects.
        
        Args:
        input_data: The input data or file path.
        
        Returns:
        A list of Document objects or an error message if the input type is unsupported.
        """
        file_extension = os.path.splitext(input_data)[1].lower() if os.path.isfile(input_data) else None
        
        try:
            if file_extension == ".pdf":
                loader = PyPDFLoader(input_data)
                docs = loader.load()
            elif file_extension == ".pptx":
                prs = PptxPresentation(input_data)
                docs = [Document(page_content="\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]), metadata={"title": slide.shapes.title.text if slide.shapes.title else ""}) for slide in prs.slides]
            elif file_extension == ".docx":
                doc = DocxDocument(input_data)
                docs = [Document(page_content=para.text, metadata={}) for para in doc.paragraphs]
            elif file_extension == ".txt":
                with open(input_data, 'r', encoding='utf-8') as file:
                    text_data = file.read()
                docs = [Document(page_content=text_data, metadata={})]
            else:
                docs = [Document(page_content=input_data, metadata={})]
            return docs
        except Exception as e:
            return f"Error loading document: {str(e)}"

    def summarize_document(docs):
        """
        Summarize a list of Document objects using a predefined LLM chain.
        
        Args:
        docs: A list of Document objects to summarize.
        
        Returns:
        The summary of the documents.
        """
        # Define the prompt template for the summary
        prompt_template = """
                Kindly provide a comprehensive and well-structured summary of the text below:
                {text}

                Use the following Markdown format:
                # Insert Descriptive Report Title

                The summary can include important subheadings

                For each subheading, use numbered bullet points (1 to 10) to ensure clarity and coherence.

                Each subheading can include sections formatted as **Section Title** if needed.

                For example:

                ## Overview
                1. ...
                2. ...
                - **Section Title**
                    1. ...
                    2. ...
                """
        
        prompt = PromptTemplate.from_template(prompt_template)
        
        # Set up the LLM chain and the document summarization chain
        llm_chain = LLMChain(llm=llm, prompt=prompt)
        stuff_chain = StuffDocumentsChain(llm_chain=llm_chain, document_variable_name="text")
        
        # Invoke the chain to get the summary
        response = stuff_chain.invoke({"input_documents": docs})
        
        return response["output_text"]

    # Define the state structure
    class GraphState(TypedDict):
        input_data: str
        input_type: str
        docs: List[Document]
        summary: str

    # Function to detect the input type and load a document
    def load_document_node(state: GraphState) -> GraphState:
        input_data = state["input_data"]
        file_extension = os.path.splitext(input_data)[1].lower() if os.path.isfile(input_data) else None
        
        try:
            if file_extension == ".pdf":
                loader = PyPDFLoader(input_data)
                state["docs"] = loader.load()
                state["input_type"] = "pdf"
            elif file_extension == ".pptx":
                prs = PptxPresentation(input_data)
                state["docs"] = [Document(page_content="\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")]), metadata={"title": slide.shapes.title.text if slide.shapes.title else ""}) for slide in prs.slides]
                state["input_type"] = "pptx"
            elif file_extension == ".docx":
                doc = DocxDocument(input_data)
                state["docs"] = [Document(page_content=para.text, metadata={}) for para in doc.paragraphs]
                state["input_type"] = "docx"
            elif file_extension == ".txt":
                with open(input_data, 'r', encoding='utf-8') as file:
                    text_data = file.read()
                state["docs"] = [Document(page_content=text_data, metadata={})]
                state["input_type"] = "txt"
            else:
                state["docs"] = [Document(page_content=input_data, metadata={})]
                state["input_type"] = "text"
        except Exception as e:
            state["docs"] = []
            print(f"Error loading document: {str(e)}")

        return state

    # Function to summarize a document
    def summarize_document_node(state: GraphState) -> GraphState:
        prompt_template = """
            Kindly provide a comprehensive and well-structured summary of the text below:
            {text}

            Use the following Markdown format:
            # Insert Descriptive Report Title

            The summary can include important subheadings

            For each subheading, use numbered bullet points (1 to 10) to ensure clarity and coherence.

            Each subheading can include sections formatted as **Section Title** if needed.

            For example:

            ## Overview
            1. ...
            2. ...
            - **Section Title**
                1. ...
                2. ...
            """
        
        prompt = PromptTemplate.from_template(prompt_template)
        llm_chain = LLMChain(llm=llm, prompt=prompt)
        stuff_chain = StuffDocumentsChain(llm_chain=llm_chain, document_variable_name="text")
        
        response = stuff_chain.invoke({"input_documents": state["docs"]})
        state["summary"] = response["output_text"]
        
        return state

    # Function to print the state
    def state_printer(state: GraphState) -> GraphState:
        print("---STATE PRINTER---")
        #print(f"Input Data: {state.get('input_data')}")
        #print(f"Input Type: {state.get('input_type')}")
        #print(f"Documents: {state.get('docs')}")
        #print(f"Summary: {state.get('summary')}")
        return state

    # Define the workflow
    workflow = StateGraph(GraphState)

    # Add nodes to the workflow
    workflow.add_node("load_document_node", load_document_node)
    workflow.add_node("summarize_document_node", summarize_document_node)
    workflow.add_node("state_printer", state_printer)

    # Set entry point
    workflow.set_entry_point("load_document_node")

    # Define edges
    workflow.add_edge("load_document_node", "summarize_document_node")
    workflow.add_edge("summarize_document_node", "state_printer")
    workflow.add_edge("state_printer", END)

    # Compile the workflow
    app = workflow.compile()

    return app